"""YouTube動画からキャプチャとテキストを抽出するモジュール。

このモジュールはYouTube動画から定期的に画像をキャプチャし、
対応する字幕とともにスライドに変換します。
"""

import os
import re
import time
from datetime import datetime

import cv2
import numpy as np
import yt_dlp
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from youtube_transcript_api import YouTubeTranscriptApi

# 追加: Whisper API連携モジュールをインポート
try:
    from whisper_integration import WhisperTranscriptionProvider, TranscriptQualitySelector
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False


class EnhancedYouTubeTutorialExtractor:
    """YouTubeチュートリアル動画から画像とテキストを抽出し、スライドを生成するクラス。"""

    def __init__(
            self,
            url,
            output_dir="output",
            interval=30,
            lang="ja",
            format_type="pptx",
            compact_text=True,
            max_text_length=200,
            screen_change_threshold=0.6,
            image_quality=95,
            video_format="bestvideo[ext=mp4]+bestaudio[ext=m4a]/best[ext=mp4]",
            add_thumbnail=True,
            use_whisper=False,
            whisper_api_key=None,
            whisper_model="medium",
            force_whisper=False):
        """初期化メソッド。

        Args:
            url: YouTube動画のURL
            output_dir: 出力ディレクトリ
            interval: フレーム抽出間隔（秒）
            lang: 字幕の言語コード
            format_type: 出力形式 ("pptx" または "slides")
            compact_text: テキストを改行せずにコンパクトに表示するか
            max_text_length: スライド分割の文字数閾値
            screen_change_threshold: 画面変化の検出閾値 (0-1の間, 高いほど敏感)
            image_quality: 画像の品質 (1-100, 高いほど高品質)
            video_format: ダウンロードする動画の形式とクオリティ
            add_thumbnail: サムネイルを追加するかどうか
            use_whisper: Whisper APIを使用するかどうか
            whisper_api_key: Whisper API（OpenAI API）のキー
            whisper_model: 使用するWhisperモデル
            force_whisper: 品質評価にかかわらず常にWhisper APIの結果を使用するか
        """
        self.url = url
        self.output_dir = output_dir
        self.interval = interval
        self.lang = lang
        self.format_type = format_type
        self.compact_text = compact_text
        self.max_text_length = max_text_length
        self.screen_change_threshold = screen_change_threshold
        self.image_quality = image_quality
        self.video_format = video_format
        self.add_thumbnail = add_thumbnail
        self.video_id = self.extract_video_id()  # メソッド名を修正
        self.temp_dir = os.path.join(output_dir, "temp")
        
        # 追加: Whisper API関連の設定
        self.use_whisper = use_whisper and WHISPER_AVAILABLE
        self.whisper_api_key = whisper_api_key
        self.whisper_model = whisper_model
        self.force_whisper = force_whisper
        
        # 必要なディレクトリを作成
        os.makedirs(self.output_dir, exist_ok=True)
        os.makedirs(self.temp_dir, exist_ok=True)
        
    def extract_video_id(self):  # メソッド名を修正
        """YouTubeのURLからビデオIDを抽出する。

        Returns:
            str: YouTube動画のID
            
        Raises:
            ValueError: URLが無効な場合
        """
        if "youtu.be" in self.url:
            video_id = self.url.split("/")[-1]
            if "?" in video_id:
                video_id = video_id.split("?")[0]
            return video_id
        elif "youtube.com" in self.url:
            match = re.search(r"v=([^&]+)", self.url)
            if match:
                return match.group(1)
            else:
                raise ValueError("YouTubeのビデオIDが見つかりませんでした")
        else:
            raise ValueError("無効なYouTube URLです")
    
    def get_video_info(self):
        """動画のタイトルと説明を取得する。

        Returns:
            dict: 動画のメタデータを含む辞書
        """
        try:
            # yt-dlpの設定オプション
            ydl_opts = {
                'skip_download': True,
                'quiet': True,
                'no_warnings': True,
            }
            
            # yt-dlpを使って情報を取得
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(self.url, download=False)
            
            # 情報を取得
            title = info.get('title', '不明なタイトル')
            author = info.get('uploader', '不明な作成者')
            description = info.get('description', '')
            length = info.get('duration', 0)
            
            # サムネイル画像のURLを取得
            thumbnails = info.get('thumbnails', [])
            thumbnail_url = None
            # 利用可能な最高品質のサムネイルを探す
            if thumbnails:
                # 解像度で並べ替え（高解像度のものを優先）
                sorted_thumbnails = sorted(
                    [t for t in thumbnails if 'width' in t and 'height' in t],
                    key=lambda x: x.get('width', 0) * x.get('height', 0),
                    reverse=True
                )
                if sorted_thumbnails:
                    thumbnail_url = sorted_thumbnails[0].get('url')
            
            # サムネイルのダウンロード
            thumbnail_path = None
            if thumbnail_url:
                try:
                    import requests
                    from PIL import Image
                    import io
                    
                    # まずメモリにダウンロード
                    response = requests.get(thumbnail_url, stream=True)
                    if response.status_code == 200:
                        # 一時的なWebPファイル
                        temp_thumbnail = os.path.join(self.temp_dir, f"{self.video_id}_thumbnail_temp")
                        with open(temp_thumbnail, 'wb') as f:
                            for chunk in response.iter_content(1024):
                                f.write(chunk)
                        
                        # PILを使って画像を開き、JPG形式で保存（PowerPoint互換）
                        thumbnail_path = os.path.join(self.temp_dir, f"{self.video_id}_thumbnail.jpg")
                        try:
                            img = Image.open(temp_thumbnail)
                            # RGBAの場合はRGBに変換（透過部分は白に）
                            if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                                background = Image.new('RGB', img.size, (255, 255, 255))
                                background.paste(img, mask=img.split()[3] if img.mode == 'RGBA' else None)
                                background.save(thumbnail_path, 'JPEG', quality=95)
                            else:
                                img.convert('RGB').save(thumbnail_path, 'JPEG', quality=95)
                            
                            # 一時ファイルを削除
                            if os.path.exists(temp_thumbnail):
                                os.remove(temp_thumbnail)
                        except Exception as e:
                            print(f"サムネイル画像の変換に失敗しました: {e}")
                            if os.path.exists(temp_thumbnail):
                                os.remove(temp_thumbnail)
                            thumbnail_path = None
                    else:
                        print(f"サムネイル取得エラー: HTTPステータス {response.status_code}")
                        thumbnail_path = None
                except Exception as e:
                    print(f"サムネイルのダウンロードに失敗しました: {e}")
                    thumbnail_path = None
            
            # 公開日の処理
            upload_date = info.get('upload_date')
            if upload_date and len(upload_date) == 8:
                try:
                    year = int(upload_date[:4])
                    month = int(upload_date[4:6])
                    day = int(upload_date[6:8])
                    publish_date = datetime(year, month, day)
                except:
                    publish_date = None
            else:
                publish_date = None
            
            return {
                "title": title,
                "author": author,
                "description": description,
                "length": length,
                "publish_date": publish_date,
                "thumbnail_path": thumbnail_path
            }
        except Exception as e:
            print(f"動画情報の取得に失敗しました: {e}")
            # 最低限の情報を返す
            return {
                "title": f"動画 {self.video_id}",
                "author": "不明",
                "description": "",
                "length": 0,
                "publish_date": None,
                "thumbnail_path": None
            }
    
    def download_transcript(self, video_path=None):
        """字幕をダウンロードする。
        
        Whisper APIが有効な場合は、YouTubeから取得した字幕とWhisper APIで生成した字幕を比較し、
        より品質の高い方を返します。force_whisperが有効な場合は常にWhisper APIの結果を優先します。

        Args:
            video_path: 動画ファイルのパス（Whisper API使用時に必要）
            
        Returns:
            list: 字幕データのリスト
        """
        youtube_transcript = []
        
        # YouTubeから字幕を取得
        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(self.video_id)
            
            # 指定された言語の字幕を探す
            transcript = None
            for t in transcript_list:
                if t.language_code == self.lang:
                    transcript = t
                    break
            
            # 指定言語が見つからない場合は自動生成字幕を探す
            if transcript is None:
                for t in transcript_list:
                    if t.is_generated:
                        transcript = t
                        if t.language_code != self.lang:
                            transcript = t.translate(self.lang)
                        break
            
            # それでも見つからなければデフォルトのものを使用
            if transcript is None:
                transcript = transcript_list.find_transcript([self.lang, 'en'])
            
            # 字幕データを取得
            transcript_data = transcript.fetch()
            
            # データ形式を確認して標準化（辞書形式に統一）
            standardized_data = []
            for entry in transcript_data:
                # 既に辞書形式かチェック
                if isinstance(entry, dict) and 'start' in entry and 'text' in entry:
                    standardized_data.append({
                        'start': float(entry['start']),
                        'text': str(entry['text']),
                        'duration': float(entry.get('duration', 0))
                    })
                # オブジェクト形式の場合
                elif hasattr(entry, 'start') and hasattr(entry, 'text'):
                    standardized_data.append({
                        'start': float(entry.start),
                        'text': str(entry.text),
                        'duration': float(getattr(entry, 'duration', 0))
                    })
                # その他の場合はスキップ
                else:
                    print(f"未対応の字幕データ形式です: {type(entry)}")
            
            youtube_transcript = standardized_data
        except Exception as e:
            print(f"YouTubeからの字幕ダウンロードに失敗しました: {e}")
        
        # Whisper APIが有効で、API KEYが設定されている場合は音声認識を実行
        whisper_transcript = []
        if self.use_whisper and self.whisper_api_key and video_path:
            try:
                print("Whisper APIで音声認識を実行しています...")
                whisper_provider = WhisperTranscriptionProvider(
                    api_key=self.whisper_api_key,
                    language=self.lang,
                    model=self.whisper_model
                )
                whisper_transcript = whisper_provider.get_transcript(video_path, self.temp_dir)
            except Exception as e:
                print(f"Whisper APIによる音声認識に失敗しました: {e}")
        
        # YouTubeとWhisperの両方から字幕が取得できた場合、選択ロジックを適用
        if youtube_transcript and whisper_transcript:
            return TranscriptQualitySelector.select_best_transcript(
                youtube_transcript, whisper_transcript, force_whisper=self.force_whisper)
        # どちらか一方しか取得できなかった場合はそれを返す
        elif youtube_transcript:
            return youtube_transcript
        elif whisper_transcript:
            return whisper_transcript
        # どちらも取得できなかった場合は空リストを返す
        else:
            return []
    
    def compute_frame_similarity(self, frame1, frame2):
        """2つのフレーム間の類似度を計算する。

        Args:
            frame1: 比較元のフレーム
            frame2: 比較先のフレーム
            
        Returns:
            float: 類似度 (0-1, 1が完全一致)
        """
        try:
            # リサイズして計算を高速化
            small1 = cv2.resize(frame1, (64, 36))
            small2 = cv2.resize(frame2, (64, 36))
            
            # グレースケールに変換
            gray1 = cv2.cvtColor(small1, cv2.COLOR_BGR2GRAY)
            gray2 = cv2.cvtColor(small2, cv2.COLOR_BGR2GRAY)
            
            # ヒストグラム比較
            hist1 = cv2.calcHist([gray1], [0], None, [256], [0, 256])
            hist2 = cv2.calcHist([gray2], [0], None, [256], [0, 256])
            
            cv2.normalize(hist1, hist1, 0, 1, cv2.NORM_MINMAX)
            cv2.normalize(hist2, hist2, 0, 1, cv2.NORM_MINMAX)
            
            # ヒストグラム類似度（コサイン距離）
            similarity = cv2.compareHist(hist1, hist2, cv2.HISTCMP_CORREL)
            
            return max(0, min(1, similarity))  # 0-1の範囲に収める
        except Exception as e:
            print(f"類似度計算エラー: {e}")
            return 1.0  # エラーの場合は変化なしとする
    
    def extract_frames_advanced(self):
        """動画からフレームを抽出し、インターバル、テキスト量、画面変化に基づいてフレームを選別する。

        Returns:
            tuple: (フレームパスのリスト, フレーム時間のリスト, 対応するテキストのリスト)
        """
        try:
            print("動画をダウンロード中...")
            video_path = os.path.join(self.temp_dir, f"{self.video_id}.mp4")
            
            # yt-dlpの設定オプション（よりロバストな設定）
            ydl_opts = {
                'format': 'best',  # デフォルトはbestフォーマット
                'outtmpl': video_path,
                'quiet': True,
                'no_warnings': True,
            }
            
            # 動画のダウンロードを試みる（複数の方法で）
            success = False
            
            # 1. 指定されたフォーマットでダウンロードを試みる
            if self.video_format != 'best':
                try:
                    format_opts = ydl_opts.copy()
                    format_opts['format'] = self.video_format
                    with yt_dlp.YoutubeDL(format_opts) as ydl:
                        ydl.download([self.url])
                        success = os.path.exists(video_path)
                        if success:
                            print("指定されたフォーマットでダウンロードしました")
                except Exception as e:
                    print(f"指定フォーマットでのダウンロードに失敗しました: {e}")
            
            # 2. 'best' フォーマットでのダウンロードを試みる
            if not success:
                try:
                    print("利用可能な最高品質のフォーマットを使用します...")
                    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                        ydl.download([self.url])
                        success = os.path.exists(video_path)
                        if success:
                            print("最高品質フォーマットでダウンロードしました")
                except Exception as e:
                    print(f"最高品質フォーマットでのダウンロードに失敗しました: {e}")
            
            # 3. フォーマット指定をさらに緩和してダウンロードを試みる
            if not success:
                try:
                    print("一般的なフォーマットでダウンロードを試みます...")
                    alt_opts = ydl_opts.copy()
                    alt_opts['format'] = 'mp4'
                    with yt_dlp.YoutubeDL(alt_opts) as ydl:
                        ydl.download([self.url])
                        success = os.path.exists(video_path)
                        if success:
                            print("mp4フォーマットでダウンロードしました")
                except Exception as e:
                    print(f"mp4フォーマットでのダウンロードに失敗しました: {e}")
            
            # 4. より簡素なフォーマット指定での最終試行
            if not success:
                try:
                    print("最終試行: シンプルなフォーマットでダウンロードします...")
                    last_opts = ydl_opts.copy()
                    last_opts['format'] = 'worstvideo+worstaudio/worst'
                    with yt_dlp.YoutubeDL(last_opts) as ydl:
                        ydl.download([self.url])
                        success = os.path.exists(video_path)
                        if success:
                            print("最低品質フォーマットでダウンロードしました")
                except Exception as e:
                    print(f"最終試行でのダウンロードにも失敗しました: {e}")
            
            if not os.path.exists(video_path):
                raise Exception("動画のダウンロードに失敗しました")
            
            print("字幕をダウンロード中...")
            transcript = self.download_transcript(video_path)
            if not transcript:
                print("警告: 字幕を取得できませんでした。テキスト分析に基づくスライド分割は無効になります。")
                
            print("フレームを抽出中...")
            cap = cv2.VideoCapture(video_path)
            fps = cap.get(cv2.CAP_PROP_FPS)
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = total_frames / fps
            
            # フレーム情報を記録するための辞書
            # キー: 時間（秒）、値: (フレームパス, テキスト)
            # 辞書を使用することで同一時間のフレームの重複を自動的に排除
            frame_dict = {}
            
            # 最初のフレームを取得
            prev_frame = None
            
            # インターバルごとにフレームを抽出
            for sec in range(0, int(duration), self.interval):
                # 現在の時間ポイントを設定
                cap.set(cv2.CAP_PROP_POS_MSEC, sec * 1000)
                ret, frame = cap.read()
                
                if not ret:
                    continue
                
                # 画面変化検出（最初のフレーム以外）
                if prev_frame is not None:
                    similarity = self.compute_frame_similarity(prev_frame, frame)
                    
                    # 類似度が閾値より低い（= 変化が大きい）場合はフレームを追加
                    if similarity < self.screen_change_threshold:
                        # 画像ファイルとして高品質で保存
                        timestamp = time.strftime('%H:%M:%S', time.gmtime(sec))
                        img_path = os.path.join(self.temp_dir, f"frame_{sec}.jpg")
                        
                        # 高画質パラメータを設定
                        img_params = [cv2.IMWRITE_JPEG_QUALITY, self.image_quality]
                        cv2.imwrite(img_path, frame, img_params)
                        
                        # テキスト取得
                        text = self._get_transcript_at_time(transcript, sec, window_size=15)
                        
                        # フレーム情報を辞書に追加（重複を避けるため）
                        if sec not in frame_dict:
                            frame_dict[sec] = (img_path, text)
                            print(f"画面変化検出: {sec}秒 (類似度: {similarity:.2f})")
                
                # 定期的なインターバルでのフレーム保存（高品質）
                # 既に同じ時間のフレームが存在する場合はスキップ
                if sec not in frame_dict:
                    timestamp = time.strftime('%H:%M:%S', time.gmtime(sec))
                    img_path = os.path.join(self.temp_dir, f"frame_{sec}.jpg")
                    
                    # 高画質パラメータを設定
                    img_params = [cv2.IMWRITE_JPEG_QUALITY, self.image_quality]
                    cv2.imwrite(img_path, frame, img_params)
                    
                    # テキスト取得
                    text = self._get_transcript_at_time(transcript, sec, window_size=15)
                    
                    # フレーム情報を辞書に追加
                    frame_dict[sec] = (img_path, text)
                
                # テキスト量に基づく分割処理
                current_time = sec
                
                # frame_dictから現在の時間に対応するテキストを取得
                current_text = frame_dict.get(sec, (None, ""))[1]
                
                if transcript and current_text and len(current_text) > self.max_text_length:
                    # テキスト量が多い場合、セグメントを分割して追加のフレームを作成
                    segments = self._split_text_by_amount(
                        transcript, current_time, current_time + self.interval)
                    
                    for seg_time in segments[1:]:  # 最初のセグメントは既に追加済み
                        # セグメント時間が既に存在する場合はスキップ
                        if seg_time in frame_dict:
                            continue
                            
                        # フレームを取得
                        cap.set(cv2.CAP_PROP_POS_MSEC, seg_time * 1000)
                        ret, seg_frame = cap.read()
                        if ret:
                            # 画像ファイルとして高品質で保存
                            timestamp = time.strftime('%H:%M:%S', time.gmtime(seg_time))
                            img_path = os.path.join(self.temp_dir, f"frame_{seg_time}.jpg")
                            
                            # 高画質パラメータを設定
                            img_params = [cv2.IMWRITE_JPEG_QUALITY, self.image_quality]
                            cv2.imwrite(img_path, seg_frame, img_params)
                            
                            # テキスト取得
                            seg_text = self._get_transcript_at_time(
                                transcript, seg_time, window_size=15)
                            
                            # フレーム情報を辞書に追加
                            frame_dict[seg_time] = (img_path, seg_text)
                            print(f"テキスト量による分割: {seg_time}秒")
                
                # 現在のフレームを保存
                prev_frame = frame
            
            cap.release()
            
            # 辞書から時間でソートされたリストを作成
            sorted_times = sorted(frame_dict.keys())
            frames = []
            frame_times = []
            transcript_chunks = []
            
            for t in sorted_times:
                img_path, text = frame_dict[t]
                frames.append(img_path)
                frame_times.append(t)
                transcript_chunks.append(text)
            
            print(f"合計 {len(frames)} フレームを抽出しました")
            return frames, frame_times, transcript_chunks
        
        except Exception as e:
            print(f"フレーム抽出に失敗しました: {e}")
            import traceback
            traceback.print_exc()  # スタックトレースを表示
            return [], [], []

    def _split_text_by_amount(self, transcript, start_time, end_time):
        """指定時間範囲内の字幕を文字量に基づいて分割し、分割点の時間を返す。

        Args:
            transcript: 字幕データ
            start_time: 開始時間（秒）
            end_time: 終了時間（秒）
            
        Returns:
            list: 分割点の時間リスト
        """
        if not transcript:
            return [start_time]
            
        # 時間範囲内の字幕エントリを抽出
        time_range_entries = [
            entry for entry in transcript 
            if start_time <= entry['start'] < end_time
        ]
        
        if not time_range_entries:
            return [start_time]
            
        # 時間順にソート
        time_range_entries.sort(key=lambda x: x['start'])
        
        # 分割点を格納するリスト（最初の時間点を含む）
        split_points = [start_time]
        
        # 文字数を累積していく
        accumulated_text = ""
        
        for entry in time_range_entries:
            accumulated_text += " " + entry['text']
            
            # 文字数が閾値を超えたら分割点として時間を追加
            if len(accumulated_text) >= self.max_text_length:
                split_points.append(entry['start'])
                accumulated_text = ""
        
        return split_points
    
    def _get_transcript_at_time(self, transcript, time_sec, window_size=30):
        """指定された時間付近の字幕テキストを取得する。

        Args:
            transcript: 字幕データ
            time_sec: 検索する時間（秒）
            window_size: 検索する時間範囲（秒）
            
        Returns:
            str: テキスト
        """
        if not transcript:
            return ""
            
        relevant_parts = []
        last_end_time = -1
        
        # 時間順に字幕をソート
        sorted_transcript = sorted(transcript, key=lambda x: x['start'])
        
        for entry in sorted_transcript:
            start_time = entry['start']
            
            # 指定時間範囲内のエントリーを探す
            if abs(start_time - time_sec) <= window_size:
                text = entry['text'].strip()
                if not text:
                    continue
                
                # コンパクトモードの場合、近い時間の字幕をまとめる
                if self.compact_text and relevant_parts and start_time - last_end_time < 2.0:
                    # 直前の字幕と結合（改行なし）
                    relevant_parts[-1] += " " + text
                else:
                    relevant_parts.append(text)
                
                last_end_time = start_time + entry.get('duration', 0)
        
        # スペースで結合（改行なし）
        return " ".join(relevant_parts) if self.compact_text else "\n".join(relevant_parts)
    
    def create_google_slides(self, frames_data, video_info):
        """GoogleSlides用プレゼンテーション（PPTX形式）を作成する。

        Args:
            frames_data: フレームデータのタプル (フレームパス, 時間, テキスト)
            video_info: 動画情報の辞書
            
        Returns:
            str: 作成したPPTXファイルのパス
        """
        try:
            frames, frame_times, transcript_chunks = frames_data
            
            prs = Presentation()
            
            # タイトルスライド - 空白レイアウトを使用して完全に手動で構成
            title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
            
            # サムネイル画像の追加（利用可能かつオプションが有効な場合）
            thumbnail_height = 0
            if self.add_thumbnail and video_info.get("thumbnail_path") and os.path.exists(video_info["thumbnail_path"]):
                try:
                    # サムネイルをタイトルスライドの上部に追加
                    thumbnail_width = prs.slide_width * 0.6  # スライド幅の60%
                    thumbnail_top = Inches(0.5)
                    thumbnail_left = (prs.slide_width - thumbnail_width) / 2  # 中央揃え
                    
                    # サムネイルの高さを計算
                    img = Image.open(video_info["thumbnail_path"])
                    img_width, img_height = img.size
                    aspect_ratio = img_height / img_width
                    thumbnail_height = thumbnail_width * aspect_ratio
                    
                    # サムネイルを追加
                    thumbnail = title_slide.shapes.add_picture(
                        video_info["thumbnail_path"], 
                        thumbnail_left, 
                        thumbnail_top, 
                        width=thumbnail_width
                    )
                    print(f"サムネイル画像を追加しました")
                except Exception as e:
                    print(f"サムネイル画像の追加に失敗しました: {e}")
                    thumbnail_height = 0
            
            # タイトルをテキストボックスとして追加
            title_top = Inches(0.5) + thumbnail_height + Inches(0.3)  # サムネイルの下に配置
            left = Inches(0.5)
            width = prs.slide_width - Inches(1.0)
            height = Inches(1.0)
            
            title_box = title_slide.shapes.add_textbox(left, title_top, width, height)
            tf = title_box.text_frame
            p = tf.add_paragraph()
            p.text = video_info["title"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.alignment = 1  # 中央揃え
            
            # サブタイトルをテキストボックスとして追加
            subtitle_top = title_top + Inches(1.0)  # タイトルの下に配置
            subtitle_box = title_slide.shapes.add_textbox(left, subtitle_top, width, height)
            tf = subtitle_box.text_frame
            p = tf.add_paragraph()
            p.text = f"作成者: {video_info['author']}"
            p.font.size = Pt(18)
            p.alignment = 1  # 中央揃え
            
            p = tf.add_paragraph()
            p.text = f"抽出日: {datetime.now().strftime('%Y-%m-%d')}"
            p.font.size = Pt(18)
            p.alignment = 1  # 中央揃え
            
            # 使用した音声認識エンジンの情報を追加
            if self.use_whisper and self.whisper_api_key:
                p = tf.add_paragraph()
                whisper_info = f"音声認識: Whisper API ({self.whisper_model}モデル)"
                if self.force_whisper:
                    whisper_info += " - 優先使用"
                p.text = whisper_info
                p.font.size = Pt(14)
                p.alignment = 1  # 中央揃え
            
            # YouTube URL
            url_top = subtitle_top + Inches(0.8)
            url_box = title_slide.shapes.add_textbox(left, url_top, width, Inches(0.5))
            tf = url_box.text_frame
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = "元の動画を見る"
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0, 0, 255)
            run.hyperlink.address = self.url
            p.alignment = 1  # 中央揃え
            
            # 目次スライド - 空白レイアウトを使用して完全に手動で構成
            toc_slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
            
            # 目次タイトルをテキストボックスとして追加
            left = Inches(0.5)
            top = Inches(0.5)
            width = prs.slide_width - Inches(1.0)
            height = Inches(0.6)
            
            title_box = toc_slide.shapes.add_textbox(left, top, width, height)
            tf = title_box.text_frame
            p = tf.add_paragraph()
            p.text = "目次"
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = 1  # 中央揃え
            
            # 目次リストをテキストボックスとして追加
            left = Inches(1.0)
            top = Inches(1.5)
            width = prs.slide_width - Inches(2.0)
            height = Inches(5.0)
            
            list_box = toc_slide.shapes.add_textbox(left, top, width, height)
            tf = list_box.text_frame
            tf.word_wrap = True
            
            # 目次項目を追加
            if frame_times and len(frame_times) > 0:
                # 最初の目次項目
                timestamp = time.strftime('%H:%M:%S', time.gmtime(frame_times[0]))
                preview = transcript_chunks[0][:30] + "..." if transcript_chunks[0] and len(transcript_chunks[0]) > 30 else transcript_chunks[0] or "..."
                
                p = tf.add_paragraph()
                p.text = f"1. [{timestamp}] {preview}"
                p.font.size = Pt(14)
                # 左インデント
                p.level = 0
                
                # 残りの目次項目を追加
                for i in range(1, len(frame_times)):
                    timestamp = time.strftime('%H:%M:%S', time.gmtime(frame_times[i]))
                    preview = transcript_chunks[i][:30] + "..." if transcript_chunks[i] and len(transcript_chunks[i]) > 30 else transcript_chunks[i] or "..."
                    
                    p = tf.add_paragraph()
                    p.text = f"{i+1}. [{timestamp}] {preview}"
                    p.font.size = Pt(14)
                    # 左インデント
                    p.level = 0
            else:
                p = tf.add_paragraph()
                p.text = "目次項目はありません"
                p.font.size = Pt(14)
            
            # コンテンツスライド - 空白レイアウトを使用
            for i, (frame_path, frame_time, transcript_text) in enumerate(
                    zip(frames, frame_times, transcript_chunks)):
                # 空白のスライドを追加
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
                timestamp = time.strftime('%H:%M:%S', time.gmtime(frame_time))
                
                # タイトルをテキストボックスとして手動で追加（上部に配置）
                left = Inches(0.5)
                top = Inches(0.1)  # 位置を上に移動
                width = prs.slide_width - Inches(1.0)
                height = Inches(0.5)
                
                title_box = slide.shapes.add_textbox(left, top, width, height)
                tf = title_box.text_frame
                p = tf.add_paragraph()
                p.text = f"セクション {i+1} - {timestamp}"
                p.font.size = Pt(20)
                p.font.bold = True
                p.alignment = 1  # 中央揃え
                
                # フレーム画像（スライドの80%のサイズ）
                img_height = Inches(3.0)  # デフォルト値
                if os.path.exists(frame_path):
                    try:
                        # スライドのサイズを取得
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height
                        
                        # スライドの画質向上のため、解像度を考慮した適切なサイズ比率を計算
                        try:
                            # 画像ファイルを読み込む
                            image = cv2.imread(frame_path)
                            if image is not None:
                                frame_height, frame_width = image.shape[:2]
                                slide_aspect_ratio = 16/9  # スライドのアスペクト比（標準的なワイドスクリーン）
                                
                                # アスペクト比を考慮して配置する画像の幅を決定
                                if frame_width/frame_height > slide_aspect_ratio:
                                    # 横長の画像の場合、幅を基準に調整（横幅の80%に拡大）
                                    img_width = int(slide_width * 0.8)
                                    img_height = int(img_width * frame_height / frame_width)
                                else:
                                    # 縦長または正方形の画像の場合、高さを基準に調整（高さの60%に拡大）
                                    img_height = int(slide_height * 0.6)
                                    img_width = int(img_height * frame_width / frame_height)
                                    # 幅が横幅の80%を超える場合は調整
                                    if img_width > int(slide_width * 0.8):
                                        img_width = int(slide_width * 0.8)
                                        img_height = int(img_width * frame_height / frame_width)
                            else:
                                # 画像が読み込めない場合はデフォルト値を使用
                                img_width = int(slide_width * 0.8)
                                img_height = int(slide_height * 0.5)
                        except Exception as e:
                            print(f"画像サイズ計算エラー: {e}")
                            # エラー時はデフォルト値を使用
                            img_width = int(slide_width * 0.8)  # 横幅の80%に拡大
                            img_height = int(slide_height * 0.6)
                        
                        # 画像を適切な位置に配置
                        img_left = int((slide_width - img_width) / 2)
                        img_top = Inches(0.7)  # タイトルとの間隔を適切に設定
                        
                        pic = slide.shapes.add_picture(frame_path, img_left, img_top, width=img_width)
                        
                        # 画像の高さを取得
                        img_height = pic.height
                    except Exception as e:
                        print(f"画像の追加に失敗しました: {e}")
                        # img_heightはデフォルト値のままにする
                        
                # トランスクリプトテキスト（画像の後ろにかからないように下部に配置）
                if transcript_text:
                    left = Inches(0.5)
                    text_top = Inches(0.8) + img_height + Inches(0.2)  # 画像の下に余白を設けて配置
                    width = prs.slide_width - Inches(1.0)
                    height = Inches(1.5)
                    
                    text_box = slide.shapes.add_textbox(left, text_top, width, height)
                    tf = text_box.text_frame
                    tf.word_wrap = True
                    p = tf.add_paragraph()
                    p.text = transcript_text
                    p.font.size = Pt(14)
                
                # YouTube リンク（スライドの最下部に配置、クリック可能なハイパーリンクとして）
                youtube_link = f"{self.url}&t={int(frame_time)}s"
                link_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
                tf = link_box.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "この部分を動画で見る"
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0, 0, 255)
                run.hyperlink.address = youtube_link
            
            # PPTXを保存（GoogleSlides用）
            safe_title = ''.join(c for c in video_info["title"] if c.isalnum() or c.isspace() or c == '_')[:30]
            pptx_path = os.path.join(self.output_dir, f"{safe_title}_googleslides.pptx")
            prs.save(pptx_path)
            print(f"GoogleSlides用ファイルを作成しました: {pptx_path}")
            return pptx_path
        
        except Exception as e:
            print(f"GoogleSlides用ファイル作成に失敗しました: {e}")
            import traceback
            traceback.print_exc()  # 詳細なエラー情報を表示
            return None
    
    def process(self):
        """メイン処理を実行する。

        Returns:
            str: 作成したファイルのパス、または失敗時はNone
        """
        print(f"YouTube URL: {self.url}の処理を開始します")
        
        # 動画情報を取得
        video_info = self.get_video_info()
        print(f"動画タイトル: {video_info['title']}")
        
        # Whisper 強制優先モードの場合は表示
        if self.use_whisper and self.force_whisper:
            print("Whisper API優先モードが有効です: YouTubeの字幕よりWhisper APIの結果を優先します")
        
        # 拡張機能付きフレーム抽出（テキスト量・画面変化による分割を含む）
        frames_data = self.extract_frames_advanced()
        if not frames_data[0]:
            print("フレームを抽出できませんでした。処理を中止します。")
            return None
        
        result_path = None
        
        # GoogleSlides形式でスライドを生成
        result_path = self.create_google_slides(frames_data, video_info)
        
        # 一時ファイルを削除
        for frame_path in frames_data[0]:
            if os.path.exists(frame_path):
                os.remove(frame_path)
        
        if os.path.exists(os.path.join(self.temp_dir, f"{self.video_id}.mp4")):
            os.remove(os.path.join(self.temp_dir, f"{self.video_id}.mp4"))
        
        return result_path