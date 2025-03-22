import os
import cv2
import argparse
import time
import yt_dlp
from youtube_transcript_api import YouTubeTranscriptApi
from PIL import Image
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re
from datetime import datetime

class YouTubeTutorialExtractor:
    def __init__(self, url, output_dir="output", interval=30, lang="ja", format_type="pptx", compact_text=True):
        """
        YouTubeチュートリアル動画から画像とテキストを抽出し、スライドファイルを生成するクラス
        
        Args:
            url (str): YouTube動画のURL
            output_dir (str): 出力ディレクトリ
            interval (int): フレーム抽出間隔（秒）
            lang (str): 字幕の言語コード
            format_type (str): 出力形式 ("pptx" または "slides")
            compact_text (bool): テキストを改行せずにコンパクトに表示するか
        """
        self.url = url
        self.output_dir = output_dir
        self.interval = interval
        self.lang = lang
        self.format_type = format_type
        self.compact_text = compact_text
        self.video_id = self._extract_video_id()
        self.temp_dir = os.path.join(output_dir, "temp")
        
        # 必要なディレクトリを作成
        os.makedirs(self.output_dir, exist_ok=True)
        os.makedirs(self.temp_dir, exist_ok=True)
        
    def _extract_video_id(self):
        """YouTubeのURLからビデオIDを抽出"""
        if "youtu.be" in self.url:
            video_id = self.url.split("/")[-1]
            if "?" in video_id:
                video_id = video_id.split("?")[0]
            return video_id
        elif "youtube.com" in self.url:
            match = re.search(r"v=([^&]+)", self.url)
            if match:
                return match.group(1)
            else:
                raise ValueError("YouTubeのビデオIDが見つかりませんでした")
        else:
            raise ValueError("無効なYouTube URLです")
    
    def get_video_info(self):
        """動画のタイトルと説明を取得（yt-dlpライブラリを使用）"""
        try:
            # yt-dlpの設定オプション
            ydl_opts = {
                'skip_download': True,
                'quiet': True,
                'no_warnings': True,
            }
            
            # yt-dlpを使って情報を取得
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(self.url, download=False)
            
            # 情報を取得
            title = info.get('title', '不明なタイトル')
            author = info.get('uploader', '不明な作成者')
            description = info.get('description', '')
            length = info.get('duration', 0)
            
            # 公開日の処理
            upload_date = info.get('upload_date')
            if upload_date and len(upload_date) == 8:
                try:
                    year = int(upload_date[:4])
                    month = int(upload_date[4:6])
                    day = int(upload_date[6:8])
                    publish_date = datetime(year, month, day)
                except:
                    publish_date = None
            else:
                publish_date = None
            
            return {
                "title": title,
                "author": author,
                "description": description,
                "length": length,
                "publish_date": publish_date
            }
        except Exception as e:
            print(f"動画情報の取得に失敗しました: {e}")
            # 最低限の情報を返す
            return {
                "title": f"動画 {self.video_id}",
                "author": "不明",
                "description": "",
                "length": 0,
                "publish_date": None
            }
    
    def download_transcript(self):
        """字幕をダウンロード"""
        try:
            transcript_list = YouTubeTranscriptApi.list_transcripts(self.video_id)
            
            # 指定された言語の字幕を探す
            transcript = None
            for t in transcript_list:
                if t.language_code == self.lang:
                    transcript = t
                    break
            
            # 指定言語が見つからない場合は自動生成字幕を探す
            if transcript is None:
                for t in transcript_list:
                    if t.is_generated:
                        transcript = t
                        if t.language_code != self.lang:
                            transcript = t.translate(self.lang)
                        break
            
            # それでも見つからなければデフォルトのものを使用
            if transcript is None:
                transcript = transcript_list.find_transcript([self.lang, 'en'])
            
            # 字幕データを取得
            transcript_data = transcript.fetch()
            
            # データ形式を確認して標準化（辞書形式に統一）
            standardized_data = []
            for entry in transcript_data:
                # 既に辞書形式かチェック
                if isinstance(entry, dict) and 'start' in entry and 'text' in entry:
                    standardized_data.append({
                        'start': float(entry['start']),
                        'text': str(entry['text']),
                        'duration': float(entry.get('duration', 0))
                    })
                # オブジェクト形式の場合
                elif hasattr(entry, 'start') and hasattr(entry, 'text'):
                    standardized_data.append({
                        'start': float(entry.start),
                        'text': str(entry.text),
                        'duration': float(getattr(entry, 'duration', 0))
                    })
                # その他の場合はスキップ
                else:
                    print(f"未対応の字幕データ形式です: {type(entry)}")
            
            return standardized_data
        except Exception as e:
            print(f"字幕のダウンロードに失敗しました: {e}")
            return []
    
    def extract_frames(self):
        """動画からフレームを抽出 (yt-dlpライブラリを使用)"""
        try:
            print("動画をダウンロード中...")
            video_path = os.path.join(self.temp_dir, f"{self.video_id}.mp4")
            
            # yt-dlpの設定オプション
            ydl_opts = {
                'format': 'best[ext=mp4]',
                'outtmpl': video_path,
                'quiet': True,
                'no_warnings': True,
            }
            
            # yt-dlpを使って動画をダウンロード
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([self.url])
            
            if not os.path.exists(video_path):
                raise Exception("動画のダウンロードに失敗しました")
                
            print("フレームを抽出中...")
            cap = cv2.VideoCapture(video_path)
            fps = cap.get(cv2.CAP_PROP_FPS)
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = total_frames / fps
            
            frames = []
            frame_times = []
            
            # インターバルごとにフレームを抽出
            for sec in range(0, int(duration), self.interval):
                cap.set(cv2.CAP_PROP_POS_MSEC, sec * 1000)
                ret, frame = cap.read()
                if ret:
                    # 画像ファイルとして保存
                    timestamp = time.strftime('%H:%M:%S', time.gmtime(sec))
                    img_path = os.path.join(self.temp_dir, f"frame_{sec}.jpg")
                    cv2.imwrite(img_path, frame)
                    
                    frames.append(img_path)
                    frame_times.append(sec)
            
            cap.release()
            return frames, frame_times
        
        except Exception as e:
            print(f"フレーム抽出に失敗しました: {e}")
            import traceback
            traceback.print_exc()  # スタックトレースを表示
            return [], []
    
    def _get_transcript_at_time(self, transcript, time_sec, window_size=30):
        """
        指定された時間付近の字幕テキストを取得
        Args:
            transcript: 字幕データ
            time_sec: 検索する時間（秒）
            window_size: 検索する時間範囲（秒）
        """
        if not transcript:
            return ""
            
        relevant_parts = []
        last_end_time = -1
        
        # 時間順に字幕をソート
        sorted_transcript = sorted(transcript, key=lambda x: x['start'])
        
        for entry in sorted_transcript:
            start_time = entry['start']
            
            # 指定時間範囲内のエントリーを探す
            if abs(start_time - time_sec) <= window_size:
                text = entry['text'].strip()
                if not text:
                    continue
                
                # コンパクトモードの場合、近い時間の字幕をまとめる
                if self.compact_text and relevant_parts and start_time - last_end_time < 2.0:
                    # 直前の字幕と結合（改行なし）
                    relevant_parts[-1] += " " + text
                else:
                    relevant_parts.append(text)
                
                last_end_time = start_time + entry.get('duration', 0)
        
        # スペースで結合（改行なし）
        return " ".join(relevant_parts) if self.compact_text else "\n".join(relevant_parts)
    
    def create_pptx(self, frames, transcript, video_info):
        """PowerPointプレゼンテーションを作成"""
        try:
            prs = Presentation()
            
            # タイトルスライド
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            title.text = video_info["title"]
            subtitle.text = f"作成者: {video_info['author']}\n" \
                           f"抽出日: {datetime.now().strftime('%Y-%m-%d')}"
            
            # 目次スライド
            toc_slide = prs.slides.add_slide(prs.slide_layouts[1])
            toc_slide.shapes.title.text = "目次"
            toc_content = toc_slide.placeholders[1]
            toc_text = ""
            
            for i, frame_time in enumerate(frames[1]):
                timestamp = time.strftime('%H:%M:%S', time.gmtime(frame_time))
                transcript_text = self._get_transcript_at_time(transcript, frame_time)
                preview = transcript_text[:30] + "..." if transcript_text else "..."
                toc_text += f"{i+1}. [{timestamp}] {preview}\n"
            
            toc_content.text = toc_text
            
            # コンテンツスライド
            for i, (frame_path, frame_time) in enumerate(zip(frames[0], frames[1])):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                timestamp = time.strftime('%H:%M:%S', time.gmtime(frame_time))
                
                # タイトル（小さめのフォントサイズ、より上部に配置）
                title_shape = slide.shapes.title
                title_shape.text = f"セクション {i+1} - {timestamp}"
                title_shape.top = Inches(0.3)  # タイトルを上部に移動
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(24)  # フォントサイズを小さくする
                
                # フレーム画像（スライドの80%のサイズ）
                if os.path.exists(frame_path):
                    try:
                        # スライドのサイズを取得
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height
                        
                        # 画像を80%のサイズで配置
                        img_width = int(slide_width * 0.8)
                        left = int((slide_width - img_width) / 2)
                        top = Inches(1.0)  # タイトルの下（上部に寄せる）
                        
                        pic = slide.shapes.add_picture(frame_path, left, top, width=img_width)
                    except Exception as e:
                        print(f"画像の追加に失敗しました: {e}")
                
                # トランスクリプトテキスト（画像の後ろにかからないよう下部に配置）
                transcript_text = self._get_transcript_at_time(transcript, frame_time)
                if transcript_text:
                    img_height = Inches(4.0)  # 画像の想定高さ
                    
                    left = Inches(0.5)
                    top = top + img_height + Inches(0.2)  # 画像の下に配置（余白を設定）
                    width = prs.slide_width - Inches(1.0)
                    height = Inches(1.5)
                    
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.add_paragraph()
                    p.text = transcript_text
                    p.font.size = Pt(14)
                
                # YouTube リンク（スライドの最下部に配置）
                youtube_link = f"{self.url}&t={int(frame_time)}s"
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = f"この部分を動画で見る: {youtube_link}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0, 0, 255)
            
            # PPTXを保存
            safe_title = ''.join(c for c in video_info["title"] if c.isalnum() or c.isspace() or c == '_')[:30]
            pptx_path = os.path.join(self.output_dir, f"{safe_title}_tutorial.pptx")
            prs.save(pptx_path)
            print(f"PowerPointファイルを作成しました: {pptx_path}")
            return pptx_path
        
        except Exception as e:
            print(f"PowerPoint作成に失敗しました: {e}")
            import traceback
            traceback.print_exc()  # 詳細なエラー情報を表示
            return None
    
    def create_google_slides(self, frames, transcript, video_info):
        """GoogleSlides用プレゼンテーション（PPTX形式）を作成"""
        try:
            prs = Presentation()
            
            # タイトルスライド（GoogleSlides互換レイアウト）
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # タイトルスライド
            
            # タイトルを手動でテキストボックスとして追加（レイアウトに依存しない）
            left = Inches(0.5)
            top = Inches(1.0)
            width = prs.slide_width - Inches(1.0)
            height = Inches(1.0)
            
            title_box = title_slide.shapes.add_textbox(left, top, width, height)
            tf = title_box.text_frame
            p = tf.add_paragraph()
            p.text = video_info["title"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.alignment = 1  # 中央揃え
            
            # サブタイトルを手動で追加
            left = Inches(0.5)
            top = Inches(2.5)
            width = prs.slide_width - Inches(1.0)
            height = Inches(1.0)
            
            subtitle_box = title_slide.shapes.add_textbox(left, top, width, height)
            tf = subtitle_box.text_frame
            p = tf.add_paragraph()
            p.text = f"作成者: {video_info['author']}"
            p.font.size = Pt(18)
            
            p = tf.add_paragraph()
            p.text = f"抽出日: {datetime.now().strftime('%Y-%m-%d')}"
            p.font.size = Pt(18)
            
            # 目次スライド
            toc_slide = prs.slides.add_slide(prs.slide_layouts[1])  # タイトルと内容レイアウト
            
            # 目次タイトルをプレースホルダに設定するか、テキストボックスで追加
            if len(toc_slide.placeholders) > 0:
                # プレースホルダがある場合
                try:
                    title_placeholder = toc_slide.shapes.title
                    title_placeholder.text = "目次"
                except:
                    # プレースホルダがなければテキストボックスで追加
                    title_box = toc_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
                    tf = title_box.text_frame
                    p = tf.add_paragraph()
                    p.text = "目次"
                    p.font.size = Pt(28)
                    p.font.bold = True
            else:
                # プレースホルダがない場合はテキストボックスで追加
                title_box = toc_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
                tf = title_box.text_frame
                p = tf.add_paragraph()
                p.text = "目次"
                p.font.size = Pt(28)
                p.font.bold = True
            
            # 目次テキストを追加
            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1.0)
            height = Inches(5.0)
            
            txBox = toc_slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            for i, frame_time in enumerate(frames[1]):
                timestamp = time.strftime('%H:%M:%S', time.gmtime(frame_time))
                transcript_text = self._get_transcript_at_time(transcript, frame_time, window_size=15)
                preview = transcript_text[:30] + "..." if transcript_text else "..."
                
                p = tf.add_paragraph()
                p.text = f"{i+1}. [{timestamp}] {preview}"
            
            # コンテンツスライド - 単純化したレイアウトを使用
            for i, (frame_path, frame_time) in enumerate(zip(frames[0], frames[1])):
                # 空白のスライドを追加
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
                timestamp = time.strftime('%H:%M:%S', time.gmtime(frame_time))
                
                # タイトルをテキストボックスとして手動で追加
                left = Inches(0.5)
                top = Inches(0.2)
                width = prs.slide_width - Inches(1.0)
                height = Inches(0.5)
                
                title_box = slide.shapes.add_textbox(left, top, width, height)
                tf = title_box.text_frame
                p = tf.add_paragraph()
                p.text = f"セクション {i+1} - {timestamp}"
                p.font.size = Pt(20)
                p.font.bold = True
                
                # フレーム画像（スライドの80%のサイズ）
                if os.path.exists(frame_path):
                    try:
                        # スライドのサイズを取得
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height
                        
                        # 画像を80%のサイズで配置
                        img_width = int(slide_width * 0.8)
                        img_left = int((slide_width - img_width) / 2)
                        img_top = Inches(0.8)  # タイトルの下
                        
                        pic = slide.shapes.add_picture(frame_path, img_left, img_top, width=img_width)
                        
                        # 画像の高さを取得
                        img_height = pic.height
                    except Exception as e:
                        print(f"画像の追加に失敗しました: {e}")
                        img_height = Inches(3.0)  # デフォルト値
                else:
                    img_height = Inches(3.0)  # 画像がない場合のデフォルト値
                
                # トランスクリプトテキスト（画像の後ろにかからないように下部に配置）
                transcript_text = self._get_transcript_at_time(transcript, frame_time, window_size=15)
                if transcript_text:
                    left = Inches(0.5)
                    text_top = Inches(0.8) + img_height + Inches(0.2)  # 画像の下に余白を設けて配置
                    width = prs.slide_width - Inches(1.0)
                    height = Inches(1.5)
                    
                    text_box = slide.shapes.add_textbox(left, text_top, width, height)
                    tf = text_box.text_frame
                    tf.word_wrap = True
                    p = tf.add_paragraph()
                    p.text = transcript_text
                    p.font.size = Pt(14)
                
                # YouTube リンク（スライドの最下部に配置）
                youtube_link = f"{self.url}&t={int(frame_time)}s"
                link_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
                tf = link_box.text_frame
                p = tf.add_paragraph()
                p.text = f"この部分を動画で見る: {youtube_link}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0, 0, 255)
            
            # PPTXを保存（GoogleSlides用）
            safe_title = ''.join(c for c in video_info["title"] if c.isalnum() or c.isspace() or c == '_')[:30]
            pptx_path = os.path.join(self.output_dir, f"{safe_title}_googleslides.pptx")
            prs.save(pptx_path)
            print(f"GoogleSlides用ファイルを作成しました: {pptx_path}")
            return pptx_path
        
        except Exception as e:
            print(f"GoogleSlides用ファイル作成に失敗しました: {e}")
            import traceback
            traceback.print_exc()  # 詳細なエラー情報を表示
            return None
    
    def process(self):
        """メイン処理"""
        print(f"YouTube URL: {self.url}の処理を開始します")
        
        # 動画情報を取得
        video_info = self.get_video_info()
        print(f"動画タイトル: {video_info['title']}")
        
        # 字幕をダウンロード
        transcript = self.download_transcript()
        if not transcript:
            print("字幕を取得できませんでした。処理を中止します。")
            return None
        
        # フレームを抽出
        frames = self.extract_frames()
        if not frames[0]:
            print("フレームを抽出できませんでした。処理を中止します。")
            return None
        
        result_path = None
        
        # 出力形式に応じてスライドを生成
        if self.format_type == "pptx":
            result_path = self.create_pptx(frames, transcript, video_info)
        elif self.format_type == "slides":
            result_path = self.create_google_slides(frames, transcript, video_info)
        
        # 一時ファイルを削除
        for frame_path in frames[0]:
            if os.path.exists(frame_path):
                os.remove(frame_path)
        
        if os.path.exists(os.path.join(self.temp_dir, f"{self.video_id}.mp4")):
            os.remove(os.path.join(self.temp_dir, f"{self.video_id}.mp4"))
        
        return result_path

def main():
    parser = argparse.ArgumentParser(description="YouTube チュートリアル動画からスライドを生成")
    parser.add_argument("url", help="YouTube 動画の URL")
    parser.add_argument("--output", "-o", default="output", help="出力ディレクトリ")
    parser.add_argument("--format", "-f", choices=["pptx", "slides"], default="slides", help="出力形式")
    parser.add_argument("--interval", "-i", type=int, default=30, help="フレーム抽出間隔（秒）")
    parser.add_argument("--lang", "-l", default="ja", help="字幕の言語コード（例: ja, en）")
    parser.add_argument("--nocompact", "-nc", action="store_false", help="テキストを改行せずにまとめない")
    
    args = parser.parse_args()
    
    extractor = YouTubeTutorialExtractor(
        url=args.url,
        output_dir=args.output,
        format_type=args.format,
        interval=args.interval,
        lang=args.lang,
        compact_text=args.nocompact
    )
    
    result_path = extractor.process()
    
    if result_path:
        print("\n処理が完了しました！")
        print(f"生成ファイル: {result_path}")
    else:
        print("処理に失敗しました。")

if __name__ == "__main__":
    main()
